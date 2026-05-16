#----------------------------此处依赖不要更改------------------------------#
import os
import shutil
from clientApp.operator.operators.operator_abs import OperatorAbs
from clientApp.utils.file_info_util import process_file_info
from clientApp.utils.file_utils import is_dir_empty, is_directory_exists
#----------------------------新增依赖卸载下方------------------------------#
import json
import time
import platform
import subprocess
from concurrent.futures import ThreadPoolExecutor, as_completed
import copy
import glob
import fitz
from clientApp.api_handler.qwen_api import qwen_vl_user_prom, encode_image

"""
功能：将PPT转换为Markdown格式
处理逻辑：1.拆分PPT为单页 2.转换为高清图片 3.调用多模态大模型识别 4.输出MD文件
"""

# ==================== 内置配置 ====================
DEFAULT_CONFIG = {
    "model_url": "http://10.238.57.34:8000/api/chinaUnicom/generalCenter/intelligenceCenter/Qwen2572BHH/v1",
    "model_vl_url": "http://10.238.57.34:8000/api/chinaUnicom/generalCenter/intelligenceCenter/Qwen25VL32BHH/v1",
    "app_id": "FbmEG3Xz6s",
    "app_secret": "4keWuYTToUz5EkbRCqjPX6kfhzkbeoRU",
    "nlpt_authorization": "Bearer sk-04581b8469404a74aac2003282b39e36",
    "vl_nlpt_authorization": "Bearer sk-8bc6504d76714a6fa50b3caae268b743",
    "scene_code": "SZ-00-0005"
}


# ==================== PPT转MD核心类 ====================
class PPT2MDConverter:
    """PPT转Markdown转换器"""

    def __init__(self, ppt_file_path, md_file_path, config=None, temp_dir=None):
        """
        初始化转换器

        Args:
            ppt_file_path: PPT文件路径（完整路径，如 /path/to/input.pptx）
            md_file_path: MD文件路径（完整路径，如 /path/to/output.md）
            config: API配置字典（可选，使用DEFAULT_CONFIG）
            temp_dir: 临时文件目录（可选，默认使用 /tmp/ppt2md_temp）
        """
        self.ppt_file_path = ppt_file_path
        self.md_file_path = md_file_path
        self.config = config or DEFAULT_CONFIG.copy()

        # 临时目录配置
        if temp_dir:
            self.temp_dir = temp_dir
        else:
            self.temp_dir = '/tmp/ppt2md_temp'

        # 为当前文件创建专属临时目录
        ppt_basename = os.path.splitext(os.path.basename(ppt_file_path))[0]
        self.work_dir = os.path.join(self.temp_dir, f"{ppt_basename}_{int(time.time())}")

        # 确保MD文件的父目录存在
        md_dir = os.path.dirname(md_file_path)
        if md_dir and not os.path.exists(md_dir):
            os.makedirs(md_dir, exist_ok=True)

        # 系统信息
        self.is_windows = platform.system() == 'Windows'
        self.is_macos = platform.system() == 'Darwin'
        self.is_linux = platform.system() == 'Linux'

    def _check_libreoffice(self):
        """检查LibreOffice是否可用"""
        try:
            if self.is_macos:
                libreoffice_paths = [
                    '/Applications/LibreOffice.app/Contents/MacOS/soffice',
                    '/usr/local/bin/soffice',
                    '/opt/homebrew/bin/soffice',
                ]
                for path in libreoffice_paths:
                    if os.path.exists(path):
                        return path
                result = subprocess.run(['which', 'soffice'], capture_output=True, text=True)
                if result.returncode == 0:
                    return result.stdout.strip()
            elif self.is_linux:
                result = subprocess.run(['which', 'soffice'], capture_output=True, text=True)
                if result.returncode == 0:
                    return result.stdout.strip()
            return None
        except Exception:
            return None

    def _split_ppt_to_single_pages(self, ppt_path, output_dir):
        """将多页PPT拆分为多个单页PPT文件"""
        try:
            from pptx import Presentation

            single_ppts_dir = os.path.join(output_dir, "single_ppts")
            os.makedirs(single_ppts_dir, exist_ok=True)

            prs = Presentation(ppt_path)
            single_ppt_files = []

            for idx, slide in enumerate(prs.slides, 1):
                new_prs = Presentation()
                new_prs.slide_width = prs.slide_width
                new_prs.slide_height = prs.slide_height

                slide_layout = new_prs.slide_layouts[6]
                new_slide = new_prs.slides.add_slide(slide_layout)

                for shape in slide.shapes:
                    el = shape.element
                    newel = copy.deepcopy(el)
                    new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

                single_ppt_path = os.path.join(single_ppts_dir, f"slide_{idx:03d}.pptx")
                new_prs.save(single_ppt_path)
                single_ppt_files.append(single_ppt_path)

            return single_ppt_files

        except Exception as e:
            return None

    def _convert_pdf_to_jpg(self, pdf_path, output_dir, page_idx):
        """转换单个PDF为JPG"""
        try:
            pdf_document = fitz.open(pdf_path)

            if len(pdf_document) > 0:
                page = pdf_document[0]
                zoom = 2.67
                mat = fitz.Matrix(zoom, zoom)
                pix = page.get_pixmap(matrix=mat)

                img_path = os.path.join(output_dir, f"slide_{page_idx:03d}.jpg")
                pix.save(img_path)
                pdf_document.close()

                return page_idx, img_path, None
            else:
                pdf_document.close()
                return page_idx, None, "PDF无页面"
        except Exception as e:
            return page_idx, None, str(e)

    def convert_ppt_to_images(self, ppt_path, output_dir):
        """将PPT转换为图片"""
        os.makedirs(output_dir, exist_ok=True)

        # 检查是否已有图片（避免重复转换）
        existing_images = [f for f in os.listdir(output_dir)
                           if f.endswith('.jpg') and f.startswith('slide_')
                           and '_mllm_' not in f]
        existing_images.sort()

        if len(existing_images) >= 1:
            image_files = [os.path.join(output_dir, f) for f in existing_images]
            return image_files

        # 拆分PPT为单页文件
        single_ppt_files = self._split_ppt_to_single_pages(ppt_path, output_dir)
        if not single_ppt_files:
            return None

        # 批量转换所有单页PPT为PDF
        libreoffice_path = self._check_libreoffice()
        if not libreoffice_path:
            return None

        cmd = [
            libreoffice_path,
            '--headless',
            '--invisible',
            '--nodefault',
            '--nolockcheck',
            '--convert-to', 'pdf',
            '--outdir', output_dir
        ] + single_ppt_files

        try:
            subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=300,
                env=dict(os.environ, HOME=os.path.expanduser('~'))
            )

            time.sleep(1)

            # 并行转换所有PDF为JPG
            pdf_files = glob.glob(os.path.join(output_dir, "slide_*.pdf"))

            if len(pdf_files) == 0:
                return None

            image_files = []
            with ThreadPoolExecutor(max_workers=10) as executor:
                futures = {}
                for pdf_file in pdf_files:
                    import re
                    match = re.search(r'slide_(\d+)', os.path.basename(pdf_file))
                    if match:
                        page_idx = int(match.group(1))
                        future = executor.submit(self._convert_pdf_to_jpg, pdf_file, output_dir, page_idx)
                        futures[future] = page_idx

                for future in as_completed(futures):
                    try:
                        result_idx, img_path, error = future.result()
                        if not error:
                            image_files.append((result_idx, img_path))
                    except:
                        pass

            # 按页码排序
            image_files.sort(key=lambda x: x[0])
            image_paths = [img_path for _, img_path in image_files]

            # 清理临时文件
            single_ppts_dir = os.path.join(output_dir, "single_ppts")
            if os.path.exists(single_ppts_dir):
                shutil.rmtree(single_ppts_dir)

            for pdf_file in pdf_files:
                try:
                    os.remove(pdf_file)
                except:
                    pass

            return image_paths

        except Exception as e:
            return None

    def call_multimodal_llm(self, prompt, image_path, max_retries=3):
        """调用多模态大模型"""
        try:
            # 使用统一的API函数编码图片
            base64_img = encode_image(image_path)

            # 重试机制
            last_error = None
            for attempt in range(max_retries):
                try:
                    # 调用统一的视觉模型API
                    success, result = qwen_vl_user_prom(
                        content=prompt,
                        config=self.config,
                        base64_image=base64_img,
                        temperature=0.01,
                        model_name='Qwen2.5_VL_32B'
                    )

                    if success:
                        return True, result

                    # 记录错误信息
                    last_error = result

                    # 如果失败且还有重试机会，等待后重试
                    if attempt < max_retries - 1:
                        time.sleep((attempt + 1) * 2)

                except Exception as e:
                    # 记录异常信息
                    last_error = str(e)
                    # 异常情况下，如果还有重试机会，等待后重试
                    if attempt < max_retries - 1:
                        time.sleep(2)

            return False, f"多次重试后仍失败: {last_error if last_error else '未知错误'}"

        except Exception as e:
            return False, f"调用多模态大模型异常: {str(e)}"

    def _process_one_page(self, img_path, idx, prompt):
        """处理单个页面"""
        try:
            if not os.path.exists(img_path):
                return idx, "文件不存在"

            # 调用多模态大模型
            status, mllm_result = self.call_multimodal_llm(prompt, img_path, max_retries=3)

            if status:
                return idx, mllm_result
            else:
                return idx, f"多模态大模型调用失败: {mllm_result}"

        except Exception as e:
            return idx, f"处理异常: {str(e)}"

    def convert(self, max_workers=5):
        """
        执行PPT到Markdown的转换

        Args:
            max_workers: 并行处理的最大工作线程数

        Returns:
            转换是否成功
        """
        try:
            # 创建工作目录
            os.makedirs(self.work_dir, exist_ok=True)

            # 转换PPT为图片
            image_files = self.convert_ppt_to_images(self.ppt_file_path, self.work_dir)

            if not image_files:
                raise Exception("PPT转图片失败")

            # 并行处理所有页面
            prompt = """
根据这个ppt图片，解析文字，输出纯文字的markdown格式，
需要逻辑清晰，不能遗漏省略任何内容，
不要输出任何其他内容，不需要总结
"""

            results = {}
            with ThreadPoolExecutor(max_workers=max_workers) as executor:
                futures = {}
                for idx, img_path in enumerate(image_files, 1):
                    future = executor.submit(self._process_one_page, img_path, idx, prompt)
                    futures[future] = idx

                for future in as_completed(futures):
                    try:
                        idx, result = future.result()
                        results[idx] = result
                    except Exception as e:
                        pass

            # 组装结果
            final_content = []
            for idx in sorted(results.keys()):
                final_content.append(f"# 第{idx}页\n\n{results[idx]}\n\n")

            # 保存结果到MD文件
            with open(self.md_file_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(final_content))

            # 清理临时目录
            if os.path.exists(self.work_dir):
                try:
                    shutil.rmtree(self.work_dir)
                except:
                    pass

            return True

        except Exception as e:
            # 清理临时目录
            if os.path.exists(self.work_dir):
                try:
                    shutil.rmtree(self.work_dir)
                except:
                    pass
            raise e


# ==================== 算子类 ====================
# 算子文件名要与类名相同!!!
class NormalizePptToMarkdown(OperatorAbs):
    """
    PPT转Markdown算子
    将PPTX文件转换为Markdown格式文件
    """

    def process(self, operator_id, source_path, sink_path, param, logger, config_dict):
        res_data = []
        if is_dir_empty(source_path) or not is_directory_exists(source_path):
            logger.warning("输入目录为空或不存在")
            return res_data

        # 创建输出目录
        os.makedirs(sink_path, exist_ok=True)

        # 从config_dict中读取配置
        api_config = self._get_config_from_dict(config_dict, logger)

        # 获取配置参数
        temp_dir = '/tmp/ppt2md_temp'
        max_workers = 5

        if param:
            if isinstance(param, dict):
                # param中的api_config可以覆盖config_dict中的配置
                if 'api_config' in param:
                    param_api_config = param.get('api_config')
                    if isinstance(param_api_config, dict):
                        api_config.update(param_api_config)
                temp_dir = param.get('temp_image_path', temp_dir)
                max_workers = param.get('max_workers', 5)

        for file_name in os.listdir(source_path):
            input_path = os.path.join(source_path, file_name)

            # 处理PPTX文件时，自动将输出后缀改为md
            if file_name.lower().endswith('.pptx'):
                base_name = os.path.splitext(file_name)[0]
                output_file_name = f"{base_name}.md"
                output_path = os.path.join(sink_path, output_file_name)
            else:
                output_path = os.path.join(sink_path, file_name)

            try:
                # 仅处理PPTX格式文件，其他文件直接复制
                if file_name.lower().endswith('.pptx'):
                    # 调用核心处理逻辑
                    self.main(input_path, output_path, logger, api_config, temp_dir, max_workers)
                    res_data.append(process_file_info(operator_id, input_path, output_path, 0))
                else:
                    # 非PPTX文件：直接复制，标记为"无法处理"（状态码1）
                    shutil.copy(input_path, output_path)
                    res_data.append(process_file_info(operator_id, input_path, output_path, 1))
            except Exception as e:
                # 异常场景：复制原文件，记录错误日志，标记为"处理失败"（状态码1）
                logger.error(f"处理算子-PPT转Markdown异常: 文件名={file_name}, 错误信息={str(e)}")
                shutil.copy(input_path, output_path)
                res_data.append(process_file_info(operator_id, input_path, output_path, 1))

        return res_data

    def _get_config_from_dict(self, config_dict, logger):
        """
        从config_dict中读取配置

        Args:
            config_dict: 配置字典
            logger: 日志对象

        Returns:
            api_config: API配置字典
        """
        # 使用DEFAULT_CONFIG作为默认值
        api_config = DEFAULT_CONFIG.copy()

        # 从config_dict中读取配置（如果存在）
        if config_dict and isinstance(config_dict, dict):
            if 'model_url' in config_dict:
                api_config['model_url'] = config_dict['model_url']
            if 'model_vl_url' in config_dict:
                api_config['model_vl_url'] = config_dict['model_vl_url']
            if 'app_id' in config_dict:
                api_config['app_id'] = config_dict['app_id']
            if 'app_secret' in config_dict:
                api_config['app_secret'] = config_dict['app_secret']
            if 'nlpt_authorization' in config_dict:
                api_config['nlpt_authorization'] = config_dict['nlpt_authorization']
            if 'vl_nlpt_authorization' in config_dict:
                api_config['vl_nlpt_authorization'] = config_dict['vl_nlpt_authorization']
            if 'scene_code' in config_dict:
                api_config['scene_code'] = config_dict['scene_code']

            logger.info(f"从config_dict读取配置: model_vl_url={api_config.get('model_vl_url', 'N/A')[:50]}...")
        else:
            logger.info("config_dict为空或格式不正确，使用DEFAULT_CONFIG默认配置")

        return api_config

    def main(self, ppt_file_path, md_file_path, logger, api_config=None, temp_dir=None, max_workers=5):
        """
        核心处理逻辑（PPT转Markdown）

        Args:
            ppt_file_path: 待处理PPTX文件路径（完整文件路径）
            md_file_path: 处理后MD文件路径（完整文件路径）
            logger: 日志对象
            api_config: API配置字典（可选）
            temp_dir: 临时目录（可选）
            max_workers: 并行处理线程数（可选）
        """
        # 使用传入的配置（如果为空则使用DEFAULT_CONFIG）
        if api_config is None:
            api_config = DEFAULT_CONFIG.copy()

        # 创建转换器
        converter = PPT2MDConverter(
            ppt_file_path=ppt_file_path,
            md_file_path=md_file_path,
            config=api_config,
            temp_dir=temp_dir
        )

        # 执行转换
        converter.convert(max_workers=max_workers)

        logger.info(f"PPT转Markdown成功: {ppt_file_path} -> {md_file_path}")

